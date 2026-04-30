#!/usr/bin/env python3
"""Build a PowerPoint with one slide per CSV row: <pdb>_case_study.png under output/ or outputs/."""

from __future__ import annotations

import argparse
import csv
from pathlib import Path
from typing import Dict, List, Optional

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt


def _repo_root() -> Path:
    return Path(__file__).resolve().parent.parent


def _build_holo_pdb_groups(rows: List[Dict[str, str]]) -> Dict[str, List[Dict[str, str]]]:
    """Insertion-order grouping; key = holo_pdb.lower() — same convention as scripts/pipeline.py."""
    holo_pdb_to_rows: Dict[str, List[Dict[str, str]]] = {}
    for row in rows:
        raw = (row.get("holo_pdb") or "").strip()
        if not raw:
            continue
        holo_pdb_to_rows.setdefault(raw.lower(), []).append(row)
    return holo_pdb_to_rows


def logical_output_dirname(
    row: Dict[str, str],
    holo_pdb_to_rows: Dict[str, List[Dict[str, str]]],
) -> Optional[str]:
    """Folder basename under output(s)/: holo_pdb or holo_pdb_<n> when duplicate PDB codes exist."""
    raw_h = (row.get("holo_pdb") or "").strip()
    apo = (row.get("apo_bmrb") or "").strip()
    holo_b = (row.get("holo_bmrb") or "").strip()
    if not raw_h:
        return None
    matches = holo_pdb_to_rows.get(raw_h.lower(), [])
    if len(matches) <= 1:
        return raw_h
    duplicate_index: Optional[int] = None
    for idx, m in enumerate(matches, start=1):
        if (m.get("apo_bmrb") or "").strip() == apo and (m.get("holo_bmrb") or "").strip() == holo_b:
            duplicate_index = idx
            break
    if duplicate_index is not None:
        return f"{raw_h}_{duplicate_index}"
    return None


def resolve_case_study_png(
    repo: Path,
    row: Dict[str, str],
    holo_pdb_to_rows: Dict[str, List[Dict[str, str]]],
) -> Path | None:
    holo_pdb = (row.get("holo_pdb") or "").strip()
    if not holo_pdb:
        return None
    logical_dir = logical_output_dirname(row, holo_pdb_to_rows)
    if logical_dir is None:
        return None
    name = f"{holo_pdb}_case_study.png"
    for sub in ("output", "outputs"):
        p = repo / sub / logical_dir / name
        if p.is_file():
            return p
    return None


def add_image_slide(prs: Presentation, png_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    im = Image.open(png_path).convert("RGB")
    iw, ih = im.size
    sw = float(prs.slide_width)
    sh = float(prs.slide_height)
    aspect_img = iw / ih
    aspect_slide = sw / sh
    if aspect_img > aspect_slide:
        width = sw
        height = sw / aspect_img
    else:
        height = sh
        width = sh * aspect_img
    left = int((sw - width) / 2)
    top = int((sh - height) / 2)
    slide.shapes.add_picture(str(png_path), left, top, width=int(width), height=int(height))


def add_missing_slide(prs: Presentation, pdb: str, logical_dir: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(2))
    tf = box.text_frame
    tf.text = f"holo_pdb: {pdb}"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = (
        f"Figure not found (expected {pdb}_case_study.png under "
        f"output/{logical_dir}/ and outputs/{logical_dir}/)."
    )
    p2.font.size = Pt(18)


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument(
        "--csv",
        type=Path,
        default=_repo_root() / "data" / "CSP_UBQ_ph0.5_temp5C.csv",
        help="CSV with holo_pdb column",
    )
    ap.add_argument(
        "--out",
        type=Path,
        default=_repo_root() / "tmp" / "slides" / "ph05_temp5C_case_study" / "output.pptx",
        help="Output .pptx path",
    )
    args = ap.parse_args()
    repo = _repo_root()

    args.out.parent.mkdir(parents=True, exist_ok=True)

    with args.csv.open(newline="", encoding="utf-8") as f:
        rows = list(csv.DictReader(f))

    holo_pdb_groups = _build_holo_pdb_groups(rows)

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    for row in rows:
        pdb = (row.get("holo_pdb") or "").strip()
        if not pdb:
            continue
        png = resolve_case_study_png(repo, row, holo_pdb_groups)
        logical = logical_output_dirname(row, holo_pdb_groups) or pdb
        if png is not None:
            add_image_slide(prs, png)
        else:
            add_missing_slide(prs, pdb, logical)

    prs.save(args.out)
    print(f"Wrote {len(prs.slides)} slides -> {args.out.resolve()}")


if __name__ == "__main__":
    main()
